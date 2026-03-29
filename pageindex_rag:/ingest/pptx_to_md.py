from pathlib import Path


def pptx_to_markdown(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to index PowerPoint files") from exc

    presentation = Presentation(path)
    lines = [f"# Presentation: {Path(path).stem}"]

    for slide_number, slide in enumerate(presentation.slides, start=1):
        title = ""
        bullet_lines = []

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue

            text_parts = []
            for paragraph in shape.text_frame.paragraphs:
                text = " ".join(run.text.strip() for run in paragraph.runs if run.text.strip()).strip()
                if text:
                    text_parts.append(text)

            if not text_parts:
                continue

            block_text = " ".join(text_parts).strip()
            if not title:
                title = block_text[:180]
            else:
                bullet_lines.extend(text_parts)

        lines.append(f"## Slide {slide_number}: {title or f'Slide {slide_number}'}")

        if bullet_lines:
            for item in bullet_lines:
                lines.append(f"- {item}")
        else:
            lines.append("_No readable slide text_")

    return "\n\n".join(lines)


def convert_file(src: str, out_dir: str):
    md = pptx_to_markdown(src)
    out = Path(out_dir) / (Path(src).stem + ".md")
    out.write_text(md, encoding="utf-8")
    return out
